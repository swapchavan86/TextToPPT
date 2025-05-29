from fastapi import FastAPI, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from pptx import Presentation
from pptx.util import Inches
from pptx.shapes.autoshape import Shape
from pydantic import BaseModel
import io
from starlette.responses import StreamingResponse
import logging

# Import the new service
from .openai_service import generate_presentation_slides_with_openai

# This file (`main.py`) defines the FastAPI application, including API routes.
# It handles incoming requests, uses services like `openai_service.py` for core logic
# (e.g., interacting with the OpenAI API), and then processes the results
# (e.g., generating a PPTX file).

# Setup logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

app = FastAPI()

# Enable CORS for frontend
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

class SlideRequest(BaseModel):
    topic: str
    tone: str = "educational"

@app.post("/generate-ppt/")
async def generate_ppt(slide_request: SlideRequest):
    try:
        # Use the new service to get slide data
        slides = generate_presentation_slides_with_openai(slide_request.topic, slide_request.tone)

        if not slides: # Handle case where service returns empty slides (e.g. if OpenAI response was valid but empty)
            logger.warning("generate_presentation_slides_with_openai returned no slides.")
            # Or, if this should be an error:
            # raise HTTPException(status_code=500, detail="Failed to generate presentation content: No slides returned.")

    except HTTPException as e:
        # Re-raise HTTPExceptions from the service layer
        raise e
    except Exception as e:
        # Catch any other unexpected errors from the service or this layer
        logger.error(f"Error generating presentation: {e}", exc_info=True)
        raise HTTPException(status_code=500, detail=f"An unexpected error occurred while generating the presentation: {str(e)}")

    prs = Presentation()
    for slide_data in slides:
        slide_layout = prs.slide_layouts[1]  # Title and Content
        slide = prs.slides.add_slide(slide_layout)

        # Set title safely
        title_shape = slide.shapes.title
        if title_shape:
            title_shape.text = slide_data.get("title", "Untitled Slide")

        # Set bullet points safely
        # In openai_service.py, the prompt asks for "content" which is a string of bullet points.
        # The old main.py was expecting "bullets" as a list.
        # We need to adjust this part to match the new structure from openai_service.py
        try:
            content_shape = slide.placeholders[1]
            if isinstance(content_shape, Shape) and content_shape.text_frame:
                content_text = slide_data.get("content", "") # "content" is a single string
                if content_text:
                    # Split the content string into bullet points
                    bullet_points = [bp.strip() for bp in content_text.split('\n') if bp.strip().startswith("-")]
                    
                    # Remove the leading hyphen and any leading/trailing whitespace from each bullet point
                    cleaned_bullet_points = [bp[1:].strip() if bp.startswith("-") else bp.strip() for bp in bullet_points]

                    if cleaned_bullet_points:
                        content_shape.text_frame.text = cleaned_bullet_points[0]
                        for bullet in cleaned_bullet_points[1:]:
                            p = content_shape.text_frame.add_paragraph()
                            p.text = bullet
                    elif content_text: # Handle cases where content might not be bullet points
                        content_shape.text_frame.text = content_text


        except Exception as e:
            logger.warning(f"Error adding content to slide: {e}")

    ppt_io = io.BytesIO()
    prs.save(ppt_io)
    ppt_io.seek(0)

    return StreamingResponse(
        ppt_io,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": "attachment; filename=generated_ppt.pptx"}
    )
