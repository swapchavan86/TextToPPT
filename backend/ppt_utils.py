import io
import logging
from pptx import Presentation
# from pptx.util import Inches # Only if Inches is explicitly used, otherwise remove.

logger = logging.getLogger(__name__)

def create_presentation_from_slides_data(slides_data: list) -> io.BytesIO:
    prs = Presentation()
    for slide_data in slides_data:
        slide_layout = prs.slide_layouts[1]  # Title and Content
        slide = prs.slides.add_slide(slide_layout)

        # Set title safely
        title_shape = slide.shapes.title
        if title_shape:
            title_shape.text = slide_data.get("title", "Untitled Slide")

        # Set bullet points safely
        try:
            content_shape = slide.placeholders[1]
            # Ensure content_shape is a valid placeholder and has text_frame
            if content_shape is not None and hasattr(content_shape, 'text_frame') and content_shape.text_frame is not None:
                bullets = slide_data.get("bullets", [])
                if bullets:
                    content_shape.text_frame.text = bullets[0]
                    for bullet in bullets[1:]:
                        p = content_shape.text_frame.add_paragraph()
                        p.text = bullet
                else:
                    # If no bullets, clear any default text in the placeholder if necessary
                    content_shape.text_frame.clear() 
            elif content_shape is not None and not hasattr(content_shape, 'text_frame'):
                 logger.warning(f"Placeholder does not have text_frame: {content_shape}")
            # else:
                # If placeholder doesn't exist or bullets are empty, nothing to add to content_shape for bullets.
                # Consider logging if a placeholder was expected but not found, if applicable.
                # logger.info(f"No content or bullets for slide: {slide_data.get('title', 'Untitled Slide')}")


        except Exception as e:
            # Log the specific slide data and error for better debugging
            logger.warning(f"Error adding content to slide titled '{slide_data.get('title', 'Untitled')}': {e}")


    ppt_io = io.BytesIO()
    prs.save(ppt_io)
    ppt_io.seek(0)
    return ppt_io
